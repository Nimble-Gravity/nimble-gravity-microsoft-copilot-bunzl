"""Generates synthetic, Bunzl-shaped knowledge-worker sample files.

Real Office-format output (.xlsx/.pptx/.docx) for Copilot lab exercises. Business-segment
NAMES are public (Bunzl's own reporting structure); every FIGURE below is invented for
training and must never be treated as real. Run: .venv/bin/python generate.py
"""
import os

import openpyxl
from openpyxl.styles import Font
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt

HERE = os.path.dirname(os.path.abspath(__file__))
OUT_DIR = os.path.normpath(os.path.join(HERE, '..', '..', 'assets', 'lab-data'))

# ── Source data (single source of truth — extend here for later phases) ─────
SEGMENTS = [
    {'name': 'Grocery & Foodservice', 'budget': 4200, 'actual': 4385},
    {'name': 'Safety',                'budget': 2100, 'actual': 1960},
    {'name': 'Cleaning & Hygiene',    'budget': 1850, 'actual': 1972},
    {'name': 'Retail',                'budget': 1500, 'actual': 1410},
    {'name': 'Healthcare',            'budget': 980,  'actual': 1055},
]


def generate_xlsx():
    wb = openpyxl.Workbook()
    sheet = wb.active
    sheet.title = 'Budget vs Actual'
    headers = ['Segment', 'Budget ($000s)', 'Actual ($000s)', 'Variance ($000s)']
    for col, header in enumerate(headers, start=1):
        cell = sheet.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True)
    for row, seg in enumerate(SEGMENTS, start=2):
        sheet.cell(row=row, column=1, value=seg['name'])
        sheet.cell(row=row, column=2, value=seg['budget'])
        sheet.cell(row=row, column=3, value=seg['actual'])
        sheet.cell(row=row, column=4, value=seg['actual'] - seg['budget'])

    notes = wb.create_sheet('Notes')
    notes['A1'] = 'All figures are synthetic — generated for Copilot training, not real Bunzl financials.'

    out_path = os.path.join(OUT_DIR, 'bunzl-quarterly-budget-review.xlsx')
    wb.save(out_path)
    print(f'Wrote {out_path}')


def generate_pptx():
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = 'Quarterly Business Review'
    title_slide.placeholders[1].text = 'Synthetic sample deck — Copilot training use only'

    agenda_slide = prs.slides.add_slide(prs.slide_layouts[1])
    agenda_slide.shapes.title.text = 'Agenda'
    body = agenda_slide.placeholders[1].text_frame
    body.text = 'Segment performance'
    for line in ['Key initiatives', 'Next steps']:
        p = body.add_paragraph()
        p.text = line

    perf_slide = prs.slides.add_slide(prs.slide_layouts[1])
    perf_slide.shapes.title.text = 'Segment Performance'
    body = perf_slide.placeholders[1].text_frame
    body.text = f"{SEGMENTS[0]['name']}: actual ${SEGMENTS[0]['actual']}k vs budget ${SEGMENTS[0]['budget']}k"
    for seg in SEGMENTS[1:]:
        p = body.add_paragraph()
        p.text = f"{seg['name']}: actual ${seg['actual']}k vs budget ${seg['budget']}k"

    next_slide = prs.slides.add_slide(prs.slide_layouts[1])
    next_slide.shapes.title.text = 'Next Steps'
    body = next_slide.placeholders[1].text_frame
    body.text = 'Review variance drivers by segment'
    p = body.add_paragraph()
    p.text = 'Confirm Q3 initiative owners'

    out_path = os.path.join(OUT_DIR, 'bunzl-business-review.pptx')
    prs.save(out_path)
    print(f'Wrote {out_path}')


def generate_docx():
    doc = Document()
    title = doc.add_heading('Team Update Memo', level=1)
    doc.add_paragraph('Synthetic sample document — Copilot training use only.').italic = True
    doc.add_paragraph(
        'This quarter, our segments delivered mixed results against budget, with '
        'Grocery & Foodservice and Cleaning & Hygiene running ahead of plan while '
        'Safety and Retail came in under. Details are in the attached budget review.'
    )
    doc.add_paragraph(
        'Next steps: each segment lead reviews their variance drivers and confirms '
        'initiative owners for next quarter.'
    )

    out_path = os.path.join(OUT_DIR, 'bunzl-team-update-memo.docx')
    doc.save(out_path)
    print(f'Wrote {out_path}')


if __name__ == '__main__':
    os.makedirs(OUT_DIR, exist_ok=True)
    generate_xlsx()
    generate_pptx()
    generate_docx()
